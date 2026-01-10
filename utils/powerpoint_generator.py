"""
PowerPoint Generator - Creates presentations from research content.
Generates professional PowerPoint presentations with slides for each section.
"""
import os
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from datetime import datetime
from typing import Dict, List, Optional
from io import BytesIO

logger = logging.getLogger(__name__)


class PowerpointGenerator:
    """
    Generates PowerPoint presentations from research content.
    Creates professional slides with title, content, and formatting.
    """
    
    # Color scheme (professional blue)
    PRIMARY_COLOR = RGBColor(0, 51, 102)  # Dark blue
    ACCENT_COLOR = RGBColor(0, 102, 204)  # Light blue
    TEXT_COLOR = RGBColor(32, 32, 32)     # Dark gray
    BACKGROUND_COLOR = RGBColor(255, 255, 255)  # White
    
    # Layout indices for pptx
    TITLE_SLIDE = 0
    CONTENT_SLIDE = 1
    BLANK_SLIDE = 6
    
    def __init__(self, theme: str = "professional"):
        """
        Initialize the PowerPoint generator.
        
        Args:
            theme: Presentation theme - "professional", "modern", "academic"
        """
        self.prs = Presentation()
        self.theme = theme
        self.slide_width = Inches(10)
        self.slide_height = Inches(7.5)
        self.prs.slide_width = self.slide_width
        self.prs.slide_height = self.slide_height
        
        # Set theme-specific styles
        self._apply_theme()
    
    def _apply_theme(self):
        """Apply theme-specific styling."""
        if self.theme == "modern":
            self.PRIMARY_COLOR = RGBColor(45, 52, 54)  # Dark gray
            self.ACCENT_COLOR = RGBColor(0, 184, 148)  # Teal
        elif self.theme == "academic":
            self.PRIMARY_COLOR = RGBColor(47, 85, 151)  # Academic blue
            self.ACCENT_COLOR = RGBColor(99, 110, 114)  # Slate
    
    def create_title_slide(self, title: str, subtitle: str = "", date: str = None, author: str = ""):
        """
        Create the title slide.
        
        Args:
            title: Main presentation title
            subtitle: Subtitle text
            date: Presentation date (defaults to current date)
            author: Author name
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.TITLE_SLIDE])
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        self._style_title(title_shape)
        
        # Set subtitle
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = ""
        text_frame = subtitle_shape.text_frame
        text_frame.clear()
        
        if subtitle:
            p = text_frame.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = self.TEXT_COLOR
            p.alignment = PP_ALIGN.CENTER
        
        # Add date and author at bottom
        if date or author:
            content = ""
            if author:
                content += author
            if date:
                if content:
                    content += " | "
                content += date
            
            bottom_box = slide.shapes.add_textbox(
                Inches(1), Inches(6.5), Inches(8), Inches(0.5)
            )
            tf = bottom_box.text_frame
            tf.text = content
            tf.paragraphs[0].font.size = Pt(14)
            tf.paragraphs[0].font.color.rgb = self.PRIMARY_COLOR
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        logging.info(f"Created title slide: {title}")
    
    def create_section_slide(self, section_title: str, section_number: int = None):
        """
        Create a section divider slide.
        
        Args:
            section_title: Title of the section
            section_number: Optional section number
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.BLANK_SLIDE])
        
        # Add section number
        if section_number:
            number_box = slide.shapes.add_textbox(
                Inches(1), Inches(1.5), Inches(8), Inches(1)
            )
            tf = number_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"Section {section_number}"
            p.font.size = Pt(18)
            p.font.color.rgb = self.ACCENT_COLOR
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(8), Inches(2)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = section_title
        p.font.size = Pt(40)
        p.font.color.rgb = self.PRIMARY_COLOR
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Add decorative line
        line = slide.shapes.add_shape(
            1, Inches(2), Inches(4.5), Inches(6), Inches(0.1)  # Horizontal line
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.ACCENT_COLOR
        line.line.fill.background()
        
        logging.info(f"Created section slide: {section_title}")
    
    def create_content_slide(self, title: str, content: str, bullets: List[str] = None, notes: str = ""):
        """
        Create a content slide with title and bullet points.
        
        Args:
            title: Slide title
            content: Main content text
            bullets: Optional list of bullet points
            notes: Speaker notes
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.CONTENT_SLIDE])
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        self._style_title(title_shape)
        
        # Set content
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        # Add main content paragraph
        if content:
            p = text_frame.paragraphs[0]
            p.text = content
            p.font.size = Pt(18)
            p.font.color.rgb = self.TEXT_COLOR
            p.space_after = Pt(12)
        
        # Add bullet points
        if bullets:
            for bullet in bullets:
                p = text_frame.add_paragraph()
                p.text = bullet
                p.font.size = Pt(16)
                p.font.color.rgb = self.TEXT_COLOR
                p.level = 0
                p.space_before = Pt(6)
        
        # Add speaker notes
        if notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes
        
        logging.info(f"Created content slide: {title}")
    
    def create_image_slide(self, title: str, image_path: str = None, caption: str = "", notes: str = ""):
        """
        Create a slide with an image.
        
        Args:
            title: Slide title
            image_path: Path to image file
            caption: Image caption
            notes: Speaker notes
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.BLANK_SLIDE])
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.color.rgb = self.PRIMARY_COLOR
        p.font.bold = True
        
        # Add image if path provided
        if image_path and os.path.exists(image_path):
            # Calculate image position to fit in slide
            img_width = Inches(8)
            img_box = slide.shapes.add_picture(image_path, Inches(1), Inches(1.3), width=img_width)
            
            # Add caption if provided
            if caption:
                caption_box = slide.shapes.add_textbox(
                    Inches(1), Inches(6), Inches(8), Inches(0.5)
                )
                tf = caption_box.text_frame
                p = tf.paragraphs[0]
                p.text = caption
                p.font.size = Pt(12)
                p.font.italic = True
                p.font.color.rgb = self.TEXT_COLOR
                p.alignment = PP_ALIGN.CENTER
        else:
            # Placeholder if no image
            placeholder = slide.shapes.add_shape(
                1, Inches(2), Inches(2), Inches(6), Inches(3.5)  # Rectangle
            )
            placeholder.text = "Image Placeholder"
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = RGBColor(220, 220, 220)
        
        # Add speaker notes
        if notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes
        
        logging.info(f"Created image slide: {title}")
    
    def create_table_slide(self, title: str, headers: List[str], rows: List[List[str]], notes: str = ""):
        """
        Create a slide with a table.
        
        Args:
            title: Slide title
            headers: Table column headers
            rows: Table data rows
            notes: Speaker notes
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.BLANK_SLIDE])
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.color.rgb = self.PRIMARY_COLOR
        p.font.bold = True
        
        # Calculate table dimensions
        num_rows = len(rows) + 1  # +1 for header
        num_cols = len(headers)
        col_width = Inches(8.5) / num_cols
        row_height = Inches(0.5)
        
        # Add table
        table = slide.shapes.add_table(
            num_rows, num_cols, Inches(0.75), Inches(1.3), Inches(8.5), Inches(0.5)
        ).table
        
        # Style header row
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.PRIMARY_COLOR
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(12)
        
        # Style data rows
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_text in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_text)
                cell.text_frame.paragraphs[0].font.size = Pt(11)
                cell.text_frame.paragraphs[0].font.color.rgb = self.TEXT_COLOR
        
        # Add speaker notes
        if notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes
        
        logging.info(f"Created table slide: {title}")
    
    def create_conclusion_slide(self, title: str, key_points: List[str], final_thoughts: str = ""):
        """
        Create a conclusion/summary slide.
        
        Args:
            title: Slide title (e.g., "Conclusion" or "Summary")
            key_points: List of key points to highlight
            final_thoughts: Optional closing thoughts
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.BLANK_SLIDE])
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.color.rgb = self.PRIMARY_COLOR
        p.font.bold = True
        
        # Add key points
        points_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.8), Inches(8), Inches(4)
        )
        tf = points_box.text_frame
        tf.word_wrap = True
        
        for i, point in enumerate(key_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(20)
            p.font.color.rgb = self.TEXT_COLOR
            p.space_before = Pt(12)
        
        # Add final thoughts
        if final_thoughts:
            thoughts_box = slide.shapes.add_textbox(
                Inches(1), Inches(5.8), Inches(8), Inches(1)
            )
            tf = thoughts_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = final_thoughts
            p.font.size = Pt(16)
            p.font.italic = True
            p.font.color.rgb = self.ACCENT_COLOR
            p.alignment = PP_ALIGN.CENTER
        
        logging.info(f"Created conclusion slide: {title}")
    
    def create_references_slide(self, title: str, references: str):
        """
        Create a references slide.
        
        Args:
            title: Slide title (e.g., "References" or "Bibliography")
            references: Formatted references text
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.BLANK_SLIDE])
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.color.rgb = self.PRIMARY_COLOR
        p.font.bold = True
        
        # Add references
        refs_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(1.5), Inches(8.5), Inches(5.5)
        )
        tf = refs_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = references
        p.font.size = Pt(12)
        p.font.color.rgb = self.TEXT_COLOR
        
        logging.info("Created references slide")
    
    def generate_presentation(self, sections_content: Dict[str, str], output_path: str, title: str = "", subtitle: str = "", author: str = ""):
        """
        Generate a complete presentation from research sections.
        
        Args:
            sections_content: Dictionary of section titles and content
            output_path: Path to save the presentation
            title: Presentation title
            subtitle: Presentation subtitle
            author: Author name
        """
        if not title:
            title = "Research Report"
        
        # Create title slide
        date = datetime.now().strftime("%B %d, %Y")
        self.create_title_slide(title, subtitle, date, author)
        
        # Create section slides for each section
        section_num = 1
        for section_title, content in sections_content.items():
            # Create section divider
            self.create_section_slide(section_title, section_num)
            
            # Create content slide
            # Truncate content for slide
            slide_content = content[:500] + "..." if len(content) > 500 else content
            self.create_content_slide(section_title, slide_content)
            
            section_num += 1
        
        # Create conclusion slide
        self.create_conclusion_slide(
            "Conclusion",
            ["Key findings summarized", "Recommendations based on research", "Future directions"],
            "Thank you for your attention"
        )
        
        # Save presentation
        self.save(output_path)
        logging.info(f"Generated presentation: {output_path}")
    
    def save(self, output_path: str):
        """
        Save the presentation to a file.
        
        Args:
            output_path: Path to save the .pptx file
        """
        # Ensure output directory exists
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        self.prs.save(output_path)
        logger.info(f"Saved presentation to: {output_path}")
    
    def save_to_bytes(self) -> BytesIO:
        """
        Save the presentation to a BytesIO object.
        
        Returns:
            BytesIO object containing the presentation
        """
        buffer = BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer
    
    def _style_title(self, title_shape):
        """Style a title shape according to theme."""
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.PRIMARY_COLOR
        title_shape.text_frame.paragraphs[0].font.bold = True
    
    def get_slide_count(self) -> int:
        """Return the number of slides."""
        return len(self.prs.slides)
